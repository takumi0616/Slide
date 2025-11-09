#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF => PPTX converter that preserves visual fidelity by rasterizing each page
and inserting as a full-slide image. No external binaries required.
Dependencies: pymupdf (fitz), pillow, python-pptx

Usage:
  python tools/pdf_to_pptx.py \
    --input "document/stage3/SITA学会発表_stage3_20251104_175700.pdf" \
    --output "document/stage3/SITA学会発表_stage3_20251104_175700.pptx" \
    --dpi 300
"""
import argparse
import os
import sys
import tempfile
import shutil
from typing import List

try:
    import fitz  # PyMuPDF
except ImportError as e:
    print("ERROR: PyMuPDF (pymupdf) is not installed. Please run: pip install pymupdf", file=sys.stderr)
    raise

try:
    from PIL import Image
except ImportError as e:
    print("ERROR: Pillow is not installed. Please run: pip install pillow", file=sys.stderr)
    raise

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print("ERROR: python-pptx is not installed. Please run: pip install python-pptx", file=sys.stderr)
    raise


def pdf_to_images(pdf_path: str, dpi: int, tmp_dir: str) -> List[str]:
    """
    Render PDF pages into PNG images using PyMuPDF.
    Returns list of image paths in order.
    """
    doc = fitz.open(pdf_path)
    images = []
    # scale factor: 72 dpi base in PDF space
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = os.path.join(tmp_dir, f"page-{i+1:04d}.png")
        pix.save(img_path)
        images.append(img_path)
    doc.close()
    return images


def images_to_pptx(images: List[str], out_path: str, fit_to_first_image_ratio: bool = True) -> None:
    """
    Create PPTX with each image as a full-bleed slide.
    If fit_to_first_image_ratio=True, slide size is set to the first image aspect ratio
    to minimize letterboxing.
    """
    if not images:
        raise RuntimeError("No images to insert into PPTX.")

    prs = Presentation()

    if fit_to_first_image_ratio:
        # Set slide size based on first image aspect ratio, with width ~13.333 inches (1920px @ 144dpi-ish)
        w_px, h_px = Image.open(images[0]).size
        slide_w_in = 13.333
        slide_h_in = slide_w_in * h_px / w_px
        prs.slide_width = Inches(slide_w_in)
        prs.slide_height = Inches(slide_h_in)
    else:
        # 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    for img_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Ensure output directory exists
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


def main():
    parser = argparse.ArgumentParser(description="Convert PDF to PPTX by rasterizing each page.")
    parser.add_argument("--input", "-i", required=True, help="Input PDF path")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path")
    parser.add_argument("--dpi", type=int, default=300, help="Rasterization DPI (default: 300)")
    parser.add_argument("--keep-temp", action="store_true", help="Keep temporary images directory")
    args = parser.parse_args()

    pdf_path = args.input
    out_path = args.output
    dpi = args.dpi

    if not os.path.isfile(pdf_path):
        print(f"ERROR: Input PDF not found: {pdf_path}", file=sys.stderr)
        sys.exit(1)

    tmp_dir = tempfile.mkdtemp(prefix="pdf2pptx_")
    try:
        print(f"[1/3] Rendering PDF pages at {dpi} dpi ...")
        images = pdf_to_images(pdf_path, dpi, tmp_dir)
        print(f"  - Rendered {len(images)} pages to images in: {tmp_dir}")

        print(f"[2/3] Building PPTX ...")
        images_to_pptx(images, out_path, fit_to_first_image_ratio=True)

        print(f"[3/3] Done. Saved PPTX: {out_path}")
    finally:
        if args.keep_temp:
            print(f"Temporary images kept at: {tmp_dir}")
        else:
            shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    main()
